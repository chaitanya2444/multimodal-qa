
from pdfminer.high_level import extract_text as pdf_extract
import docx
from pptx import Presentation

def extract_from_pdf(path: str) -> str:
    return pdf_extract(path)

def extract_from_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_from_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and getattr(shape, "text"):
                slides.append(shape.text)
    return "\n".join(slides)

def extract_from_txt(path: str) -> str:
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()
